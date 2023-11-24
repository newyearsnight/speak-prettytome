import os
import json
from pptx import Presentation


class PowerPointParser:

    def __init__(self, directory):
        self.directory = directory

    def parse_files(self):
        data = []
        for filename in os.listdir(self.directory):
            if filename.endswith('.pptx'):
                presentation = Presentation(os.path.join(self.directory, filename))
                slides = presentation.slides
                data.append({
                    'filename': filename,
                    'slides': []
                })

                for slide in slides:
                    data[-1]['slides'].append({
                        'title': slide.title.text,
                        'shapes': []
                    })

                    for shape in slide.shapes:
                        data[-1]['slides'][-1]['shapes'].append({
                            'name': shape.name,
                            'text': shape.text if shape.has_text else None
                        })

        return data

    def save_data(self, data, filename):
        with open(filename, 'w') as f:
            json.dump(data, f, indent=4)

if __name__ == '__main__':
    parser = PowerPointParser('powerpoints')
    data = parser.parse_files()
    parser.save_data(data, 'data.json')
