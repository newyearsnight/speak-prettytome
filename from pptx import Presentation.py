from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Hello, World!"
content.text = "Welcome to Python-pptx tutorial."

prs.save('test.pptx')